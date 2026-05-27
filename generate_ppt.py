import os
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. 讀取與處理數據
csv_path = r"c:\Users\user\Desktop\gjg\活頁簿1.csv"
df = pd.read_csv(csv_path, header=None, names=["Quarter", "Metric_A", "Metric_B"])

# 原始數據為反向時間軸（2026 Q1 -> 2015 Q1），我們將其反轉為正向時間軸以利趨勢分析
df_chrono = df.iloc[::-1].reset_index(drop=True)

# 2. 計算關鍵統計指標
avg_a = df_chrono["Metric_A"].mean()
max_a = df_chrono["Metric_A"].max()
max_a_q = df_chrono.loc[df_chrono["Metric_A"].idxmax(), "Quarter"]
min_a = df_chrono["Metric_A"].min()
min_a_q = df_chrono.loc[df_chrono["Metric_A"].idxmin(), "Quarter"]  # 2016 Q4 異常值

avg_b = df_chrono["Metric_B"].mean()
max_b = df_chrono["Metric_B"].max()
max_b_q = df_chrono.loc[df_chrono["Metric_B"].idxmax(), "Quarter"]
min_b = df_chrono["Metric_B"].min()
min_b_q = df_chrono.loc[df_chrono["Metric_B"].idxmin(), "Quarter"]

latest_q = df.iloc[0]["Quarter"]
latest_a = df.iloc[0]["Metric_A"]
latest_b = df.iloc[0]["Metric_B"]

# 3. 繪製精美的趨勢圖表
plt.style.use('dark_background')
fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(12, 6.5), sharex=True)
fig.patch.set_facecolor('#0B0F19')
ax1.set_facecolor('#0B0F19')
ax2.set_facecolor('#0B0F19')

# X 軸每隔 4 個季度顯示一個點，避免擁擠
x_ticks = df_chrono.index
x_labels = [df_chrono.loc[i, "Quarter"] if i % 4 == 0 or i == len(df_chrono)-1 else "" for i in x_ticks]

# 繪製指標 A 趨勢
ax1.plot(df_chrono.index, df_chrono["Metric_A"], color='#00F2FE', linewidth=2.5, marker='o', markersize=4, label='Metric A')
ax1.set_title("Metric A Historical Trend (2015 - 2026)", color='#E2E8F0', fontsize=12, pad=10)
ax1.grid(True, color='#1E293B', linestyle='--', alpha=0.5)
ax1.tick_params(colors='#94A3B8', labelsize=9)
for spine in ax1.spines.values():
    spine.set_color('#24324F')

# 繪製指標 B 趨勢
ax2.plot(df_chrono.index, df_chrono["Metric_B"], color='#C084FC', linewidth=2.5, marker='s', markersize=4, label='Metric B')
ax2.set_title("Metric B Historical Trend (2015 - 2026)", color='#E2E8F0', fontsize=12, pad=10)
ax2.grid(True, color='#1E293B', linestyle='--', alpha=0.5)
ax2.tick_params(colors='#94A3B8', labelsize=9)
plt.xticks(x_ticks, x_labels, rotation=45, color='#94A3B8')
for spine in ax2.spines.values():
    spine.set_color('#24324F')

plt.tight_layout()
chart_path = r"c:\Users\user\Desktop\gjg\trend_chart.png"
plt.savefig(chart_path, dpi=300, facecolor='#0B0F19')
plt.close()

# 4. 初始化 PPTX 簡報 (16:9 寬螢幕)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 色彩定義 (Dark Theme System)
BG_COLOR = RGBColor(11, 15, 25)        # 深藍黑 #0B0F19
CARD_COLOR = RGBColor(22, 29, 48)      # 卡片深灰藍 #161D30
BORDER_COLOR = RGBColor(36, 50, 79)    # 邊框灰藍 #24324F
TEXT_WHITE = RGBColor(255, 255, 255)   # 白色 #FFFFFF
TEXT_MUTED = RGBColor(148, 163, 184)   # 灰藍 #94A3B8
COLOR_CYAN = RGBColor(0, 242, 254)     # 亮青色 #00F2FE
COLOR_VIOLET = RGBColor(192, 132, 252) # 亮紫色 #C084FC

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def create_title(slide, text, subtitle_text=None):
    # 新增標題區塊
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft JhengHei"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = "Microsoft JhengHei"
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(6)

# ==================== SLIDE 1: 封面投影片 ====================
slide_layout = prs.slide_layouts[6] # 空白版面
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1)

# 背景裝飾光效 (漸層卡片感)
deco_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
deco_shape.fill.solid()
deco_shape.fill.fore_color.rgb = COLOR_CYAN
deco_shape.line.fill.background()

title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.5), Inches(3.5))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

p1 = tf.paragraphs[0]
p1.text = "季度數據分析與趨勢洞察報告"
p1.font.name = "Microsoft JhengHei"
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = TEXT_WHITE

p2 = tf.add_paragraph()
p2.text = f"基於歷史數據指標的深度解析報告 (2015 Q1 - {latest_q})"
p2.font.name = "Microsoft JhengHei"
p2.font.size = Pt(18)
p2.font.color.rgb = COLOR_CYAN
p2.space_before = Pt(16)

p3 = tf.add_paragraph()
p3.text = "數據來源: 活頁簿1.csv | 報告產出：AI 數據分析助理"
p3.font.name = "Microsoft JhengHei"
p3.font.size = Pt(12)
p3.font.color.rgb = TEXT_MUTED
p3.space_before = Pt(64)

# ==================== SLIDE 2: 執行摘要 (數據卡片) ====================
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2)
create_title(slide2, "數據執行摘要", "歷史季度數據的核心指標與概覽")

# 指標 A 卡片
card_a = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
card_a.fill.solid()
card_a.fill.fore_color.rgb = CARD_COLOR
card_a.line.color.rgb = BORDER_COLOR
card_a.line.width = Pt(1.5)

tf_a = card_a.text_frame
tf_a.word_wrap = True
tf_a.margin_left = tf_a.margin_right = tf_a.margin_top = tf_a.margin_bottom = Inches(0.3)

p = tf_a.paragraphs[0]
p.text = "指標 A 核心概覽"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_CYAN

# 最新數值
p = tf_a.add_paragraph()
p.text = f"\n最新數值 ({latest_q})"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p = tf_a.add_paragraph()
p.text = f"{latest_a:.2f}"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE

# 統計數據
p = tf_a.add_paragraph()
p.text = f"• 歷史平均值：{avg_a:.2f}\n" \
         f"• 歷史最大值：{max_a:.2f} ({max_a_q})\n" \
         f"• 歷史最小值：{min_a:.2f} ({min_a_q}) *異常值"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_WHITE
p.space_before = Pt(12)

# 指標 B 卡片
card_b = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
card_b.fill.solid()
card_b.fill.fore_color.rgb = CARD_COLOR
card_b.line.color.rgb = BORDER_COLOR
card_b.line.width = Pt(1.5)

tf_b = card_b.text_frame
tf_b.word_wrap = True
tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Inches(0.3)

p = tf_b.paragraphs[0]
p.text = "指標 B 核心概覽"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_VIOLET

# 最新數值
p = tf_b.add_paragraph()
p.text = f"\n最新數值 ({latest_q})"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p = tf_b.add_paragraph()
p.text = f"{latest_b:.2f}"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE

# 統計數據
p = tf_b.add_paragraph()
p.text = f"• 歷史平均值：{avg_b:.2f}\n" \
         f"• 歷史最大值：{max_b:.2f} ({max_b_q})\n" \
         f"• 歷史最小值：{min_b:.2f} ({min_b_q})"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_WHITE
p.space_before = Pt(12)

# ==================== SLIDE 3: 歷史趨勢圖表 ====================
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3)
create_title(slide3, "數據趨勢視覺化", "歷史走勢對比與波動分析")

# 插入圖表
slide3.shapes.add_picture(chart_path, Inches(0.8), Inches(1.8), width=Inches(8.5), height=Inches(4.8))

# 右側分析說明卡片
desc_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.6), Inches(1.8), Inches(2.9), Inches(4.8))
desc_box.fill.solid()
desc_box.fill.fore_color.rgb = CARD_COLOR
desc_box.line.color.rgb = BORDER_COLOR
desc_box.line.width = Pt(1)

tf_desc = desc_box.text_frame
tf_desc.word_wrap = True
tf_desc.margin_left = tf_desc.margin_right = tf_desc.margin_top = tf_desc.margin_bottom = Inches(0.2)

p = tf_desc.paragraphs[0]
p.text = "走勢分析要點"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE

p = tf_desc.add_paragraph()
p.text = "\n1. 指標 A 成長軌跡：\n" \
         "指標 A 整體呈現穩健上升趨勢。在 2015-2019 年維持在 50-60 水準，2020 年起攀升至 70-80 以上，近期 (2026 Q1) 達到 87.58 的高檔區間。\n\n" \
         "2. 指標 B 的週期波動：\n" \
         "指標 B 則呈現顯著的週期性與高峰。在 2021-2022 年間達到 >300 的歷史高峰，隨後在 2024-2026 年逐漸回落，目前 2026 Q1 為 71.9，處於歷史相對低點。"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(11)
p.font.color.rgb = TEXT_MUTED
p.space_before = Pt(8)

# ==================== SLIDE 4: 關鍵洞察與異常值分析 ====================
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4)
create_title(slide4, "關鍵洞察與異常分析", "聚焦數據的波動轉折與異常表現")

# 左側：2016 Q4 異常值解析
insight_1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
insight_1.fill.solid()
insight_1.fill.fore_color.rgb = CARD_COLOR
insight_1.line.color.rgb = BORDER_COLOR
insight_1.line.width = Pt(1)

tf_in1 = insight_1.text_frame
tf_in1.word_wrap = True
tf_in1.margin_left = tf_in1.margin_right = tf_in1.margin_top = tf_in1.margin_bottom = Inches(0.3)

p = tf_in1.paragraphs[0]
p.text = "⚠️ 指標 A 異常值分析 (2016 Q4)"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(239, 68, 68) # 紅色警告

p = tf_in1.add_paragraph()
p.text = f"\n• 異常數值：-0.13\n" \
         f"• 背景對比：前一季度 (2016 Q3) 為 52.48，後一季度 (2017 Q1) 為 53.06。\n\n" \
         f"• 潛在原因：\n" \
         f"   1. 資料登錄錯誤：該季度數值可能遺失或輸入負號。\n" \
         f"   2. 外部重大衝擊：該期間企業可能發生極端的單季虧損或銷量歸零事件。\n\n" \
         f"• 處理建議：\n" \
         f"   在進行時間序列建模或統計預測時，應將此點視為「極端離群值 (Outlier)」進行平滑化或剔除處理，以免扭曲整體趨勢。"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_WHITE
p.space_before = Pt(8)

# 右側：指標 B 歷史高峰與退潮
insight_2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
insight_2.fill.solid()
insight_2.fill.fore_color.rgb = CARD_COLOR
insight_2.line.color.rgb = BORDER_COLOR
insight_2.line.width = Pt(1)

tf_in2 = insight_2.text_frame
tf_in2.word_wrap = True
tf_in2.margin_left = tf_in2.margin_right = tf_in2.margin_top = tf_in2.margin_bottom = Inches(0.3)

p = tf_in2.paragraphs[0]
p.text = "📈 指標 B 週期性巨幅波動"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_VIOLET

p = tf_in2.add_paragraph()
p.text = f"\n• 歷史高峰：345.00 ({max_b_q})\n" \
         f"• 谷底走勢：近期已滑落至 71.90 ({latest_q})。\n\n" \
         f"• 現象解讀：\n" \
         f"   指標 B 在 2021 年期間出現了爆發性成長（可能受惠於疫情期間的宅經濟、供應鏈補貨潮或特殊市場溢價），並維持了約 4 個季度的超高水準。\n\n" \
         f"• 近期警訊：\n" \
         f"   自 2023 年起指標 B 呈現穩定下行軌道，2026 Q1 更創下 2020 年以來的新低，顯示該市場/指標的需求或動能已進入顯著的收縮與築底期。"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_WHITE
p.space_before = Pt(8)

# ==================== SLIDE 5: 近期季度數據對照表 ====================
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5)
create_title(slide5, "近期季度數據明細", "近兩年 (2024 - 2026) 數據對照表")

# 新增表格
rows, cols = 10, 3
left = Inches(1.5)
top = Inches(1.8)
width = Inches(10.33)
height = Inches(4.8)

table_shape = slide5.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# 設定欄寬
table.columns[0].width = Inches(3.44)
table.columns[1].width = Inches(3.44)
table.columns[2].width = Inches(3.44)

# 表頭
headers = ["季度 (Quarter)", "指標 A (Metric A)", "指標 B (Metric B)"]
for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_COLOR
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft JhengHei"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

# 填入近 9 季的數據 (df 的前 9 筆，因為原始 df 第一筆是最新)
recent_df = df.head(9)
for row_idx, row in recent_df.iterrows():
    table_row = row_idx + 1
    data_values = [row["Quarter"], f"{row['Metric_A']:.2f}", f"{row['Metric_B']:.2f}"]
    
    for col_idx, val in enumerate(data_values):
        cell = table.cell(table_row, col_idx)
        cell.fill.solid()
        # 隔行變色效果
        if table_row % 2 == 0:
            cell.fill.fore_color.rgb = CARD_COLOR
        else:
            cell.fill.fore_color.rgb = RGBColor(16, 22, 38)
            
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE

# ==================== SLIDE 6: 結論與建議 ====================
slide6 = prs.slides.add_slide(slide_layout)
set_slide_background(slide6)
create_title(slide6, "策略結論與建議", "基於趨勢的未來展望與行動指南")

conclusion_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
conclusion_card.fill.solid()
conclusion_card.fill.fore_color.rgb = CARD_COLOR
conclusion_card.line.color.rgb = BORDER_COLOR
conclusion_card.line.width = Pt(1.5)

tf_con = conclusion_card.text_frame
tf_con.word_wrap = True
tf_con.margin_left = tf_con.margin_right = tf_con.margin_top = tf_con.margin_bottom = Inches(0.4)

p = tf_con.paragraphs[0]
p.text = "🎯 關鍵策略結論"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLOR_CYAN

p = tf_con.add_paragraph()
p.text = "\n1. 指標 A 表現強勁，需防範高檔震盪：" \
         "\n   指標 A 從 2015 年的 52.83 穩步攀升至目前的 87.58。這顯示其核心業務或基本面長期向好。建議持續增資或擴大此業務線，但因目前處於歷史高位，需注意短期回檔修正的風險。"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_WHITE
p.space_before = Pt(10)

p = tf_con.add_paragraph()
p.text = "\n2. 指標 B 面臨週期底部，應著重「控產保價」與築底規劃：" \
         "\n   指標 B 從最高峰 345 下滑至目前的 71.9。這通常是市場供過於求或外部紅利退去後的硬著陸現象。短期內不宜盲目擴張或削價競爭，應精簡庫存與人力，等待供需關係重新平衡。"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_WHITE
p.space_before = Pt(10)

p = tf_con.add_paragraph()
p.text = "\n3. 數據品質稽核與治理：" \
         "\n   針對 2016 Q4 的指標 A 異常值（-0.13），建議追溯當季原始財務/銷售傳票，確認是否為公式錯誤或人工登錄失誤，並在未來的數據收集流程中加入閾值警示機制（例如：設定指標 A 下限為 30），以保障數據治理品質。"
p.font.name = "Microsoft JhengHei"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_WHITE
p.space_before = Pt(10)

# 儲存簡報
ppt_path = r"c:\Users\user\Desktop\gjg\數據分析報告.pptx"
prs.save(ppt_path)
print("PPT generated successfully!")
